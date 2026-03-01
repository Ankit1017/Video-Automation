from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from main_app.services.pptx_export.layout_planner import plan_slide_layout
from main_app.services.pptx_export.models import LayoutPlan, PptxTemplateStyle
from main_app.services.pptx_export.text_utils import normalize_text, trim_code_for_slide


class PptxDeckBuilder:
    def __init__(self, *, style: PptxTemplateStyle) -> None:
        self._style = style

    def build(self, *, topic: str, slides: list[dict[str, Any]]) -> tuple[bytes | None, str | None]:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.dml.color import RGBColor  # type: ignore
            from pptx.enum.shapes import MSO_SHAPE  # type: ignore
            from pptx.util import Inches, Pt  # type: ignore
        except Exception:
            return None, "python-pptx is not installed. Install dependencies to enable PPTX export."

        prs = Presentation()
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        self._add_title_slide(
            prs=prs,
            blank_layout=blank_layout,
            topic=topic,
            Inches=Inches,
            Pt=Pt,
            RGBColor=RGBColor,
            MSO_SHAPE=MSO_SHAPE,
        )

        for item in slides:
            plan = plan_slide_layout(slide=item if isinstance(item, dict) else {})
            self._add_content_slide(
                prs=prs,
                blank_layout=blank_layout,
                plan=plan,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                MSO_SHAPE=MSO_SHAPE,
            )

        output = BytesIO()
        prs.save(output)
        return output.getvalue(), None

    @staticmethod
    def _rgb(rgb_cls: Any, color: tuple[int, int, int]) -> Any:
        return rgb_cls(int(color[0]), int(color[1]), int(color[2]))

    def _layout(self, key: str, fallback: tuple[float, float, float, float]) -> tuple[float, float, float, float]:
        candidate = self._style.layout_presets.get(key)
        if (
            isinstance(candidate, tuple)
            and len(candidate) == 4
            and all(isinstance(item, (int, float)) for item in candidate)
        ):
            return float(candidate[0]), float(candidate[1]), float(candidate[2]), float(candidate[3])
        return fallback

    def _add_title_slide(
        self,
        *,
        prs: Any,
        blank_layout: Any,
        topic: str,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        slide = prs.slides.add_slide(blank_layout)
        self._set_slide_background(slide=slide, RGBColor=RGBColor)
        self._add_accent_bar(
            slide=slide,
            width=prs.slide_width,
            Inches=Inches,
            RGBColor=RGBColor,
            MSO_SHAPE=MSO_SHAPE,
        )

        title_x, title_y, title_w, title_h = self._layout("title_box", (0.8, 1.2, 11.8, 2.0))
        title_shape = slide.shapes.add_textbox(Inches(title_x), Inches(title_y), Inches(title_w), Inches(title_h))
        title_frame = title_shape.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = normalize_text(topic) or "Slide Deck"
        title_font = title_paragraph.font
        title_font.bold = True
        title_font.size = Pt(self._fit_title_font_size(title_paragraph.text))
        title_font.name = self._style.typography.title_font
        title_font.color.rgb = self._rgb(RGBColor, self._style.title_color)

        subtitle_shape = slide.shapes.add_textbox(Inches(0.8), Inches(3.45), Inches(11.8), Inches(0.9))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.clear()
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = self._style.branding.footer_text
        subtitle_font = subtitle_paragraph.font
        subtitle_font.size = Pt(self._style.typography.subtitle_size)
        subtitle_font.name = self._style.typography.subtitle_font
        subtitle_font.color.rgb = self._rgb(RGBColor, self._style.body_color)

        if self._style.branding.show_title_logo:
            logo_x, logo_y, logo_w, logo_h = self._layout("title_logo_box", (10.4, 0.35, 1.8, 0.7))
            self._add_logo_picture(
                slide=slide,
                logo_path=self._style.branding.logo_path,
                left=Inches(logo_x),
                top=Inches(logo_y),
                width=Inches(logo_w),
                height=Inches(logo_h),
            )

    def _add_content_slide(
        self,
        *,
        prs: Any,
        blank_layout: Any,
        plan: LayoutPlan,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        slide = prs.slides.add_slide(blank_layout)
        self._set_slide_background(slide=slide, RGBColor=RGBColor)
        self._add_accent_bar(
            slide=slide,
            width=prs.slide_width,
            Inches=Inches,
            RGBColor=RGBColor,
            MSO_SHAPE=MSO_SHAPE,
        )

        if plan.section:
            self._add_section_chip(
                slide=slide,
                section=plan.section,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                MSO_SHAPE=MSO_SHAPE,
            )

        title_x, title_y, title_w, title_h = self._layout("title_box", (0.7, 0.7, 11.9, 1.0))
        title_shape = slide.shapes.add_textbox(Inches(title_x), Inches(title_y), Inches(title_w), Inches(title_h))
        title_frame = title_shape.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = plan.title
        title_font = title_paragraph.font
        title_font.bold = True
        title_font.size = Pt(self._fit_title_font_size(plan.title))
        title_font.name = self._style.typography.title_font
        title_font.color.rgb = self._rgb(RGBColor, self._style.title_color)

        self._render_layout_body(
            slide=slide,
            plan=plan,
            Inches=Inches,
            Pt=Pt,
            RGBColor=RGBColor,
            MSO_SHAPE=MSO_SHAPE,
        )
        self._add_footer_mark(
            slide=slide,
            Inches=Inches,
            Pt=Pt,
            RGBColor=RGBColor,
        )

        if plan.speaker_notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.clear()
            notes_frame.text = plan.speaker_notes

    def _render_layout_body(
        self,
        *,
        slide: Any,
        plan: LayoutPlan,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        if plan.layout_type == "split_code" and plan.code_snippet:
            bx, by, bw, bh = self._layout("split_bullet_box", (0.7, 1.9, 5.9, 4.9))
            bullet_box = slide.shapes.add_textbox(Inches(bx), Inches(by), Inches(bw), Inches(bh))
            self._write_bullets(
                text_frame=bullet_box.text_frame,
                bullets=plan.bullets,
                preferred_font_size=self._style.typography.body_size,
                Pt=Pt,
                RGBColor=RGBColor,
            )
            self._add_code_panel(
                slide=slide,
                code_snippet=plan.code_snippet,
                code_language=plan.code_language,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                MSO_SHAPE=MSO_SHAPE,
            )
            return

        if plan.layout_type == "dual_column":
            self._render_two_column_body(
                slide=slide,
                left_title=plan.left_title or "Left",
                left_items=plan.left_items or ["No content provided."],
                right_title=plan.right_title or "Right",
                right_items=plan.right_items or ["No content provided."],
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
            )
            return

        if plan.layout_type == "timeline":
            self._render_timeline_body(
                slide=slide,
                events=plan.events,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                MSO_SHAPE=MSO_SHAPE,
            )
            return

        if plan.layout_type == "process_flow":
            self._render_process_flow_body(
                slide=slide,
                steps=plan.steps,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                MSO_SHAPE=MSO_SHAPE,
            )
            return

        if plan.layout_type == "metric_cards":
            self._render_metric_cards_body(
                slide=slide,
                cards=plan.cards,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                MSO_SHAPE=MSO_SHAPE,
            )
            return

        x, y, w, h = self._layout("body_box", (0.8, 1.9, 11.4, 4.9))
        bullet_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        self._write_bullets(
            text_frame=bullet_box.text_frame,
            bullets=plan.bullets,
            preferred_font_size=self._style.typography.body_size + 2,
            Pt=Pt,
            RGBColor=RGBColor,
        )

    def _set_slide_background(self, *, slide: Any, RGBColor: Any) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(RGBColor, self._style.background_color)

    def _add_accent_bar(
        self,
        *,
        slide: Any,
        width: int,
        Inches: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        bar_height = max(0.08, float(self._style.shapes.accent_bar_height_in))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(bar_height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(RGBColor, self._style.accent_color)
        bar.line.fill.background()

    def _add_section_chip(
        self,
        *,
        slide: Any,
        section: str,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        cx, cy, cw, ch = self._layout("section_chip_box", (9.4, 0.55, 3.0, 0.45))
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
        chip.fill.solid()
        chip.fill.fore_color.rgb = self._rgb(RGBColor, self._style.section_chip_color)
        chip.line.fill.background()

        chip_frame = chip.text_frame
        chip_frame.clear()
        chip_paragraph = chip_frame.paragraphs[0]
        chip_paragraph.text = section[:36]
        chip_font = chip_paragraph.font
        chip_font.bold = True
        chip_font.size = Pt(max(10, self._style.typography.caption_size + 1))
        chip_font.name = self._style.typography.caption_font
        chip_font.color.rgb = self._rgb(RGBColor, self._style.section_chip_text_color)

    def _add_footer_mark(
        self,
        *,
        slide: Any,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
    ) -> None:
        tx, ty, tw, th = self._layout("footer_text_box", (0.72, 6.98, 8.95, 0.26))
        footer_shape = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
        footer_frame = footer_shape.text_frame
        footer_frame.clear()
        footer_paragraph = footer_frame.paragraphs[0]
        footer_paragraph.text = self._style.branding.footer_text
        footer_font = footer_paragraph.font
        footer_font.size = Pt(self._style.typography.caption_size)
        footer_font.name = self._style.typography.caption_font
        footer_font.color.rgb = self._rgb(RGBColor, self._style.body_color)

        if self._style.branding.show_footer_logo:
            lx, ly, lw, lh = self._layout("footer_logo_box", (10.95, 6.96, 1.0, 0.3))
            self._add_logo_picture(
                slide=slide,
                logo_path=self._style.branding.footer_logo_path or self._style.branding.logo_path,
                left=Inches(lx),
                top=Inches(ly),
                width=Inches(lw),
                height=Inches(lh),
            )

    @staticmethod
    def _add_logo_picture(
        *,
        slide: Any,
        logo_path: str,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
    ) -> None:
        path = Path(str(logo_path).strip())
        if not str(path).strip():
            return
        if not path.is_file():
            return
        try:
            slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        except (AttributeError, OSError, RuntimeError, TypeError, ValueError):
            return

    def _write_bullets(
        self,
        *,
        text_frame: Any,
        bullets: list[str],
        preferred_font_size: int,
        Pt: Any,
        RGBColor: Any,
    ) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        limited_bullets = bullets[:6]
        dynamic_font = self._fit_bullet_font_size(
            preferred_size=preferred_font_size,
            bullets=limited_bullets,
        )
        for idx, bullet in enumerate(limited_bullets):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = f"- {self._truncate_line(text=bullet, max_chars=170)}"
            paragraph.level = 0
            paragraph.space_after = Pt(8)
            paragraph.line_spacing = 1.2
            paragraph_font = paragraph.font
            paragraph_font.size = Pt(dynamic_font)
            paragraph_font.name = self._style.typography.body_font
            paragraph_font.color.rgb = self._rgb(RGBColor, self._style.body_color)

    def _render_two_column_body(
        self,
        *,
        slide: Any,
        left_title: str,
        left_items: list[str],
        right_title: str,
        right_items: list[str],
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
    ) -> None:
        lhx, lhy, lhw, lhh = self._layout("left_col_heading", (0.8, 1.86, 5.4, 0.45))
        rhx, rhy, rhw, rhh = self._layout("right_col_heading", (6.2, 1.86, 5.4, 0.45))
        left_heading = slide.shapes.add_textbox(Inches(lhx), Inches(lhy), Inches(lhw), Inches(lhh))
        right_heading = slide.shapes.add_textbox(Inches(rhx), Inches(rhy), Inches(rhw), Inches(rhh))

        left_heading_frame = left_heading.text_frame
        left_heading_frame.clear()
        left_title_paragraph = left_heading_frame.paragraphs[0]
        left_title_paragraph.text = normalize_text(left_title) or "Left"
        left_title_font = left_title_paragraph.font
        left_title_font.bold = True
        left_title_font.size = Pt(max(14, self._style.typography.body_size))
        left_title_font.name = self._style.typography.body_font
        left_title_font.color.rgb = self._rgb(RGBColor, self._style.title_color)

        right_heading_frame = right_heading.text_frame
        right_heading_frame.clear()
        right_title_paragraph = right_heading_frame.paragraphs[0]
        right_title_paragraph.text = normalize_text(right_title) or "Right"
        right_title_font = right_title_paragraph.font
        right_title_font.bold = True
        right_title_font.size = Pt(max(14, self._style.typography.body_size))
        right_title_font.name = self._style.typography.body_font
        right_title_font.color.rgb = self._rgb(RGBColor, self._style.title_color)

        lbx, lby, lbw, lbh = self._layout("left_col_box", (0.8, 2.28, 5.2, 4.6))
        rbx, rby, rbw, rbh = self._layout("right_col_box", (6.2, 2.28, 5.2, 4.6))
        left_box = slide.shapes.add_textbox(Inches(lbx), Inches(lby), Inches(lbw), Inches(lbh))
        right_box = slide.shapes.add_textbox(Inches(rbx), Inches(rby), Inches(rbw), Inches(rbh))
        self._write_bullets(
            text_frame=left_box.text_frame,
            bullets=left_items[:4] or ["No content provided."],
            preferred_font_size=max(16, self._style.typography.body_size - 1),
            Pt=Pt,
            RGBColor=RGBColor,
        )
        self._write_bullets(
            text_frame=right_box.text_frame,
            bullets=right_items[:4] or ["No content provided."],
            preferred_font_size=max(16, self._style.typography.body_size - 1),
            Pt=Pt,
            RGBColor=RGBColor,
        )

    def _render_timeline_body(
        self,
        *,
        slide: Any,
        events: list[dict[str, str]],
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        rows = events[:5] if events else [{"label": "Milestone", "detail": "No timeline events provided."}]
        start_y = 2.05
        step = 0.93
        for idx, event in enumerate(rows):
            y = start_y + idx * step
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.11), Inches(0.2), Inches(0.2))
            marker.fill.solid()
            marker.fill.fore_color.rgb = self._rgb(RGBColor, self._style.accent_color)
            marker.line.fill.background()

            event_box = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(10.8), Inches(0.88))
            frame = event_box.text_frame
            frame.clear()
            header = frame.paragraphs[0]
            header.text = event.get("label", f"Milestone {idx + 1}") or f"Milestone {idx + 1}"
            header_font = header.font
            header_font.bold = True
            header_font.size = Pt(max(14, self._style.typography.body_size - 2))
            header_font.name = self._style.typography.body_font
            header_font.color.rgb = self._rgb(RGBColor, self._style.title_color)

            detail_text = event.get("detail", "")
            if detail_text:
                detail_paragraph = frame.add_paragraph()
                detail_paragraph.text = self._truncate_line(text=detail_text, max_chars=130)
                detail_paragraph.level = 0
                detail_paragraph.space_after = Pt(0)
                detail_font = detail_paragraph.font
                detail_font.size = Pt(max(12, self._style.typography.body_size - 4))
                detail_font.name = self._style.typography.body_font
                detail_font.color.rgb = self._rgb(RGBColor, self._style.body_color)

    def _render_process_flow_body(
        self,
        *,
        slide: Any,
        steps: list[dict[str, str]],
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        rows = steps[:5] if steps else [{"title": "Step 1", "detail": "No process steps provided."}]
        start_y = 1.95
        height = 0.9
        for idx, step in enumerate(rows):
            y = start_y + idx * 0.98
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(10.9), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = self._rgb(RGBColor, self._style.section_chip_color)
            box.line.fill.background()
            frame = box.text_frame
            frame.clear()
            title_paragraph = frame.paragraphs[0]
            title_paragraph.text = f"{idx + 1}. {step.get('title', 'Step')}"
            title_font = title_paragraph.font
            title_font.bold = True
            title_font.size = Pt(max(12, self._style.typography.body_size - 3))
            title_font.name = self._style.typography.body_font
            title_font.color.rgb = self._rgb(RGBColor, self._style.section_chip_text_color)
            detail = step.get("detail", "")
            if detail:
                detail_paragraph = frame.add_paragraph()
                detail_paragraph.text = self._truncate_line(text=detail, max_chars=140)
                detail_paragraph.level = 0
                detail_font = detail_paragraph.font
                detail_font.size = Pt(max(11, self._style.typography.body_size - 6))
                detail_font.name = self._style.typography.body_font
                detail_font.color.rgb = self._rgb(RGBColor, self._style.section_chip_text_color)

    def _render_metric_cards_body(
        self,
        *,
        slide: Any,
        cards: list[dict[str, str]],
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        rows = cards[:4] if cards else [{"label": "Metric", "value": "No metric cards provided.", "context": ""}]
        positions = [
            (0.8, 2.0),
            (6.3, 2.0),
            (0.8, 4.25),
            (6.3, 4.25),
        ]
        for idx, card in enumerate(rows[:4]):
            x, y = positions[idx]
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.0), Inches(2.0))
            box.fill.solid()
            box.fill.fore_color.rgb = self._rgb(RGBColor, self._style.section_chip_color)
            box.line.fill.background()
            frame = box.text_frame
            frame.clear()
            label_paragraph = frame.paragraphs[0]
            label_paragraph.text = card.get("label", "Metric")
            label_font = label_paragraph.font
            label_font.size = Pt(max(11, self._style.typography.caption_size + 1))
            label_font.bold = True
            label_font.name = self._style.typography.caption_font
            label_font.color.rgb = self._rgb(RGBColor, self._style.section_chip_text_color)
            value_paragraph = frame.add_paragraph()
            value_paragraph.text = self._truncate_line(text=card.get("value", ""), max_chars=42)
            value_font = value_paragraph.font
            value_font.size = Pt(max(14, self._style.typography.body_size))
            value_font.bold = True
            value_font.name = self._style.typography.body_font
            value_font.color.rgb = self._rgb(RGBColor, self._style.section_chip_text_color)
            context = card.get("context", "")
            if context:
                context_paragraph = frame.add_paragraph()
                context_paragraph.text = self._truncate_line(text=context, max_chars=52)
                context_font = context_paragraph.font
                context_font.size = Pt(max(10, self._style.typography.caption_size))
                context_font.name = self._style.typography.caption_font
                context_font.color.rgb = self._rgb(RGBColor, self._style.section_chip_text_color)

    def _add_code_panel(
        self,
        *,
        slide: Any,
        code_snippet: str,
        code_language: str,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
        MSO_SHAPE: Any,
    ) -> None:
        panel_left, panel_top, panel_width, panel_height = self._layout("split_code_box", (6.75, 1.72, 5.72, 5.25))
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(panel_left),
            Inches(panel_top),
            Inches(panel_width),
            Inches(panel_height),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(RGBColor, self._style.code_panel_color)
        panel.line.fill.background()

        code_lines = trim_code_for_slide(code_snippet)
        max_line_length = max((len(line) for line in code_lines), default=0)
        code_font_size = (
            self._style.typography.code_size
            if len(code_lines) <= 12 and max_line_length <= 74
            else max(10, self._style.typography.code_size - 1)
        )

        header_shape = slide.shapes.add_textbox(
            Inches(panel_left + 0.2),
            Inches(panel_top + 0.12),
            Inches(panel_width - 0.35),
            Inches(0.32),
        )
        header_frame = header_shape.text_frame
        header_frame.clear()
        header_frame.word_wrap = False
        header = header_frame.paragraphs[0]
        header.text = f"Code ({code_language.lower()})"
        header_font = header.font
        header_font.bold = True
        header_font.size = Pt(max(11, self._style.typography.caption_size + 2))
        header_font.name = self._style.typography.caption_font
        header_font.color.rgb = self._rgb(RGBColor, self._style.code_text_color)

        code_shape = slide.shapes.add_textbox(
            Inches(panel_left + 0.2),
            Inches(panel_top + 0.55),
            Inches(panel_width - 0.35),
            Inches(panel_height - 0.74),
        )
        code_frame = code_shape.text_frame
        code_frame.clear()
        code_frame.word_wrap = False
        code_frame.margin_left = 0
        code_frame.margin_right = 0
        code_frame.margin_top = 0
        code_frame.margin_bottom = 0

        for idx, line in enumerate(code_lines):
            paragraph = code_frame.paragraphs[0] if idx == 0 else code_frame.add_paragraph()
            paragraph.text = self._truncate_line(text=line, max_chars=84)
            paragraph.level = 0
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 1.0
            code_font = paragraph.font
            code_font.name = self._style.typography.code_font
            code_font.size = Pt(code_font_size)
            code_font.color.rgb = self._rgb(RGBColor, self._style.code_text_color)

    @staticmethod
    def _truncate_line(*, text: str, max_chars: int) -> str:
        value = str(text)
        if len(value) <= max_chars:
            return value
        return value[: max(3, max_chars - 3)].rstrip() + "..."

    def _fit_title_font_size(self, title: str) -> int:
        base = max(24, int(self._style.typography.title_size))
        length = len(str(title))
        if length > 130:
            return max(24, base - 8)
        if length > 100:
            return max(26, base - 6)
        if length > 80:
            return max(28, base - 4)
        if length > 60:
            return max(30, base - 2)
        return base

    @staticmethod
    def _fit_bullet_font_size(*, preferred_size: int, bullets: list[str]) -> int:
        size = max(14, int(preferred_size))
        if len(bullets) >= 5:
            size -= 2
        longest = max((len(item) for item in bullets), default=0)
        if longest > 110:
            size -= 2
        if longest > 150:
            size -= 2
        return max(12, size)

